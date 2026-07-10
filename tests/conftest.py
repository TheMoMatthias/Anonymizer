from __future__ import annotations

from pathlib import Path

import fitz
import openpyxl
import pytest
import yaml
from docx import Document
from pptx import Presentation

from anonymizer.config import DEFAULT_CONFIG_PATH
from anonymizer.engine import build_analyzer


@pytest.fixture(scope="session")
def base_config() -> dict:
    """Loads the shipped default config without touching the real
    %LOCALAPPDATA%\\Anonymizer\\config.yaml, so tests never depend on (or
    corrupt) a user's live settings. Session-scoped: no test mutates it."""
    return yaml.safe_load(DEFAULT_CONFIG_PATH.read_text(encoding="utf-8"))


@pytest.fixture(scope="session")
def analyzer(base_config):
    """Builds the Presidio analyzer once per test session (spaCy model
    loading is slow) and reuses it across all tests."""
    return build_analyzer(base_config)


@pytest.fixture
def mapping_db_path(tmp_path) -> Path:
    return tmp_path / "test_mappings.db"


@pytest.fixture
def sample_docx(tmp_path) -> Path:
    doc = Document()
    p = doc.add_paragraph()
    run = p.add_run("Hans Mueller")
    run.bold = True
    p.add_run(" ist Kunde. IBAN DE89370400440532013000.")
    doc.add_paragraph("Steuer-ID: 12 345 678 901")
    doc.sections[0].header.paragraphs[0].text = "Vertraulich - Hans Mueller"
    path = tmp_path / "sample.docx"
    doc.save(path)
    return path


@pytest.fixture
def sample_xlsx(tmp_path) -> Path:
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Main"
    ws["A1"] = "Hans Mueller"
    ws["A2"] = "DE89370400440532013000"
    hidden_ws = wb.create_sheet("Hidden")
    hidden_ws["A1"] = "Hans Mueller"
    hidden_ws.sheet_state = "hidden"
    path = tmp_path / "sample.xlsx"
    wb.save(path)
    return path


@pytest.fixture
def sample_pptx(tmp_path) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Kunde: Hans Mueller"
    slide.placeholders[1].text_frame.text = "IBAN DE89370400440532013000"
    path = tmp_path / "sample.pptx"
    prs.save(path)
    return path


@pytest.fixture
def sample_pdf(tmp_path) -> Path:
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Kunde: Hans Mueller")
    page.insert_text((72, 100), "IBAN DE89370400440532013000")
    path = tmp_path / "sample.pdf"
    doc.save(path)
    doc.close()
    return path
