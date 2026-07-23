import zipfile

from pptx import Presentation

from anonymizer.formats import pptx_handler
from anonymizer.pipeline import _output_text_blob, apply_document, scan_document

_MODERN_COMMENT = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<p188:cmLst xmlns:p188="http://schemas.microsoft.com/office/powerpoint/2018/8/main" '
    'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    "<p188:cm><p188:txBody><a:bodyPr/>"
    "<a:p><a:r><a:t>Bitte pruefen: Hans Mueller, IBAN DE89370400440532013000</a:t></a:r></a:p>"
    "</p188:txBody></p188:cm></p188:cmLst>"
)


def _add_modern_comment(path):
    with zipfile.ZipFile(path, "a") as zf:
        zf.writestr("ppt/comments/comment1.xml", _MODERN_COMMENT)


def test_detects_person_and_iban(sample_pptx, analyzer, base_config):
    grouped = scan_document(sample_pptx, analyzer, base_config).all_actionable()
    entity_types = {g.entity_type for g in grouped}
    assert "IBAN_CODE" in entity_types
    assert any(g.entity_type == "PERSON" for g in grouped)


def test_apply_replaces_slide_text(sample_pptx, analyzer, base_config, mapping_db_path):
    grouped = scan_document(sample_pptx, analyzer, base_config).all_actionable()
    for g in grouped:
        g.action = "pseudonymize"
    out_path, report_path = apply_document(sample_pptx, grouped, analyzer, base_config, mapping_db_path)

    assert out_path.exists()
    assert report_path.exists()
    prs = Presentation(out_path)
    title_text = prs.slides[0].shapes.title.text
    assert "Hans Mueller" not in title_text


def test_scans_modern_threaded_comments(sample_pptx, analyzer, base_config):
    _add_modern_comment(sample_pptx)
    units = pptx_handler.extract_text_units(sample_pptx)
    assert any("Hans Mueller" in u.text for u in units), "modern threaded comment text must be scanned"
    grouped = scan_document(sample_pptx, analyzer, base_config).all_actionable()
    assert any(g.entity_type == "IBAN_CODE" for g in grouped)


def test_apply_redacts_modern_threaded_comments(sample_pptx, analyzer, base_config, mapping_db_path):
    _add_modern_comment(sample_pptx)
    grouped = scan_document(sample_pptx, analyzer, base_config).all_actionable()
    for g in grouped:
        g.action = "anonymize"
    out_path, _ = apply_document(sample_pptx, grouped, analyzer, base_config, mapping_db_path)

    # The safety property is that NO comment PII survives anywhere in the output
    # (whether the part was redacted or dropped). apply_document's output re-scan
    # -- which now also reads modern a:t comment text -- would fail loud on any
    # residual, so reaching here at all means it verified clean.
    with zipfile.ZipFile(out_path, "r") as zf:
        blob = b"".join(zf.read(n) for n in zf.namelist())
    assert b"Hans Mueller" not in blob
    assert b"DE89370400440532013000" not in blob


def test_pptx_line_break_does_not_misalign_redaction(tmp_path, analyzer, base_config, mapping_db_path):
    """Regression (CORRUPTION/LEAK): detection ran on p.text (which includes the
    a:br vertical-tab char) but apply mapped offsets onto p.runs (a:r only), so a
    value after a soft line break was redacted at the wrong offset. Detection now
    uses the same run list; a value after a break must still be fully removed."""
    from pptx import Presentation as _P
    from pptx.oxml.ns import qn
    from pptx.util import Inches

    prs = _P()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tf = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(2)).text_frame
    p = tf.paragraphs[0]
    p.add_run().text = "Notiz:"
    p._p.append(p._p.makeelement(qn("a:br"), {}))  # soft line break between the runs
    p.add_run().text = "IBAN DE89370400440532013000"
    path = tmp_path / "break.pptx"
    prs.save(path)

    grouped = scan_document(path, analyzer, base_config).all_actionable()
    assert any("DE89370400440532013000" in g.value for g in grouped), "IBAN after break not scanned"
    for g in grouped:
        g.action = "anonymize"
    # If offsets were misaligned this would leave residual and the output re-scan
    # would fail loud; reaching a written file with the IBAN gone proves alignment.
    out_path, _ = apply_document(path, grouped, analyzer, base_config, mapping_db_path)
    out = Presentation(out_path)
    body = " ".join(
        run.text for sh in out.slides[0].shapes if sh.has_text_frame
        for para in sh.text_frame.paragraphs for run in para.runs
    )
    assert "DE89370400440532013000" not in body


def test_pptx_chart_text_does_not_survive(tmp_path, analyzer, base_config, mapping_db_path):
    """FALSE-CLEAN (B1c): a chart's title and its cached category labels live in
    ppt/charts/chart1.xml (and the embedded source workbook), which python-pptx's
    shape walk never reaches. The value was therefore never scanned, never
    removed and never re-checked -- apply reported the file verified while the
    IBAN was still sitting in the chart."""
    from pptx import Presentation as _P
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    iban = "DE89370400440532013000"
    prs = _P()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    data = CategoryChartData()
    data.categories = [f"Konto {iban}", "Sammelkonto"]
    data.add_series("Umsatz", (1.0, 2.0))
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(6), Inches(4), data
    )
    chart = graphic_frame.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = f"Depot {iban}"
    path = tmp_path / "chart.pptx"
    prs.save(path)

    grouped = scan_document(path, analyzer, base_config).all_actionable()
    assert any(iban in g.value.replace(" ", "") for g in grouped), "chart text must be surfaced for review"
    for g in grouped:
        g.action = "anonymize"
    out_path, _ = apply_document(path, grouped, analyzer, base_config, mapping_db_path)

    with zipfile.ZipFile(out_path) as zf:
        leaking = [n for n in zf.namelist() if iban.encode() in zf.read(n)]
    assert leaking == [], f"IBAN survived in chart parts: {leaking}"
    # The chart's source data is a WHOLE OOXML package nested inside a compressed
    # part, so the byte check above cannot see into it -- _output_text_blob now
    # descends, which is what makes the backstop able to fail loud on it.
    assert iban not in _output_text_blob(out_path), "IBAN survived in the embedded chart workbook"


def test_pptx_para_run_text_preserves_line_break_boundary():
    """Regression (LEAK): concatenating runs with no separator merged a name split
    across an a:br ('Klaus'<br>'Mueller' -> 'KlausMueller' -> tagged ORGANIZATION).
    The a:br must contribute a boundary to the detection text."""
    from pptx import Presentation as _P
    from pptx.oxml.ns import qn
    from pptx.util import Inches

    from anonymizer.formats.pptx_handler import _para_run_text

    prs = _P()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    p = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(2)).text_frame.paragraphs[0]
    p.add_run().text = "Klaus"
    p._p.append(p._p.makeelement(qn("a:br"), {}))
    p.add_run().text = "Mueller"
    assert _para_run_text(p) == "Klaus\nMueller"
