from __future__ import annotations

import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation

from ..core import detect_unit
from ..models import TextUnit
from pptx.oxml.ns import qn

from .run_replace import XmlRunAdapter, apply_findings_to_runs, runs_text_and_spans


class _BreakRun:
    """A no-op run standing in for an a:br soft line break: contributes "\\n" to the
    detection text so a name split across the break ("Klaus"<br>"Mueller") is not
    concatenated into one false token ("KlausMueller" -> tagged ORGANIZATION, a
    leak), but its text setter is a no-op so apply never tries to edit the break."""

    @property
    def text(self) -> str:
        return "\n"

    @text.setter
    def text(self, value: str) -> None:  # a:br has no editable text
        pass


def _para_runs_with_breaks(p) -> list:
    """The paragraph's a:r runs in document order, with a sentinel inserted for each
    a:br. Detection AND apply use this ONE list, so offsets map to the real runs
    while the line-break boundary is preserved for detection -- fixing both the
    offset drift (which p.text caused) and the word-merge (which bare p.runs caused)."""
    runs_iter = iter(p.runs)
    out: list = []
    for child in p._p:
        if child.tag == qn("a:r"):
            r = next(runs_iter, None)
            if r is not None:
                out.append(r)
        elif child.tag == qn("a:br"):
            out.append(_BreakRun())
    return out


def _para_run_text(p) -> str:
    return runs_text_and_spans(_para_runs_with_breaks(p))[0]

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
P = f"{{{P_NS}}}"
# drawingml -- modern threaded comments (PowerPoint 2018+) store their text in
# <a:t> runs, unlike the legacy <p:text> element. Scanning only p:text was the
# documented gap that let a threaded comment leak PII into an "anonymized" file.
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
A = f"{{{A_NS}}}"
COMMENT_TEXT_TAGS = (f"{P}text", f"{A}t")
COMMENT_PARTS_GLOB = "ppt/comments/"

EXTENSIONS = (".pptx", ".ppt")


def _iter_comment_text_elements(tree):
    """Yields every text-bearing element in a comment part -- both the legacy
    <p:text> and the modern threaded-comment <a:t> runs."""
    for elem in tree.iter():
        if elem.tag in COMMENT_TEXT_TAGS:
            yield elem


def _iter_text_frame_paragraphs(shape):
    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        yield from shape.text_frame.paragraphs
    if getattr(shape, "has_table", False) and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                yield from cell.text_frame.paragraphs
    if getattr(shape, "shape_type", None) is not None and hasattr(shape, "shapes"):
        for sub in shape.shapes:
            yield from _iter_text_frame_paragraphs(sub)


def _iter_paragraphs(prs: Presentation):
    for slide in prs.slides:
        for shape in slide.shapes:
            yield from _iter_text_frame_paragraphs(shape)
        if slide.has_notes_slide:
            yield from slide.notes_slide.notes_text_frame.paragraphs


def extract_text_units(path: Path) -> list[TextUnit]:
    prs = Presentation(path)
    units = []
    for i, p in enumerate(_iter_paragraphs(prs)):
        text = _para_run_text(p)
        if text.strip():
            units.append(TextUnit(id=f"p{i}", text=text))
    units.extend(_extract_comment_units(path))
    return units


def _comment_text_elements(path: Path):
    with zipfile.ZipFile(path, "r") as zf:
        names = [n for n in zf.namelist() if n.startswith(COMMENT_PARTS_GLOB)]
        for name in names:
            tree = etree.fromstring(zf.read(name))
            for text_elem in _iter_comment_text_elements(tree):
                yield name, tree, text_elem


def _extract_comment_units(path: Path) -> list[TextUnit]:
    units = []
    for i, (name, _tree, text_elem) in enumerate(_comment_text_elements(path)):
        if text_elem.text and text_elem.text.strip():
            units.append(TextUnit(id=f"comment:{name}:{i}", text=text_elem.text))
    return units


def scan(path: Path, analyzer, config) -> list:
    findings = []
    for unit in extract_text_units(path):
        findings.extend(detect_unit(analyzer, unit, config))
    return findings


def apply(path: Path, out_path: Path, decisions: dict, analyzer, config, mapping_store) -> None:
    prs = Presentation(path)
    for p in _iter_paragraphs(prs):
        runs = _para_runs_with_breaks(p)
        text = runs_text_and_spans(runs)[0]
        if not text.strip():
            continue
        findings = detect_unit(analyzer, TextUnit(id="tmp", text=text), config)
        apply_findings_to_runs(runs, findings, decisions, mapping_store)
    prs.save(out_path)
    _apply_comments(out_path, analyzer, config, decisions, mapping_store)


def _apply_comments(path: Path, analyzer, config, decisions: dict, mapping_store) -> None:
    with zipfile.ZipFile(path, "r") as zf:
        names = zf.namelist()
        contents = {n: zf.read(n) for n in names}

    any_changed = False
    for name in names:
        if not name.startswith(COMMENT_PARTS_GLOB):
            continue
        tree = etree.fromstring(contents[name])
        part_changed = False
        for text_elem in _iter_comment_text_elements(tree):
            if not text_elem.text or not text_elem.text.strip():
                continue
            unit = TextUnit(id="tmp", text=text_elem.text)
            findings = detect_unit(analyzer, unit, config)
            if not findings:
                continue
            apply_findings_to_runs([XmlRunAdapter(text_elem)], findings, decisions, mapping_store)
            part_changed = True
        if part_changed:
            contents[name] = etree.tostring(tree, xml_declaration=True, encoding="UTF-8", standalone=True)
            any_changed = True

    if any_changed:
        with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
            for name in names:
                zf.writestr(name, contents[name])
